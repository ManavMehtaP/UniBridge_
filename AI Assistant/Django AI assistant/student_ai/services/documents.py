from __future__ import annotations

import hashlib
import os
import tempfile
from pathlib import Path
from urllib.parse import urlparse

import requests

TEXT_EXTENSIONS = {".txt", ".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".c", ".cpp", ".md", ".json", ".html", ".css"}


def _is_remote(source: str | Path) -> bool:
    if isinstance(source, Path):
        return False
    parsed = urlparse(str(source))
    return parsed.scheme in {"http", "https"}


def file_hash(source: str | Path) -> str:
    if _is_remote(source):
        response = requests.get(str(source), timeout=60)
        response.raise_for_status()
        return hashlib.sha256(response.content).hexdigest()
    return hashlib.sha256(Path(source).read_bytes()).hexdigest()


def _extract_from_path(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in TEXT_EXTENSIONS:
        return path.read_text(encoding="utf-8", errors="ignore")
    if suffix == ".pdf":
        # PyMuPDF handles modern PDFs more reliably than PyPDF2 and keeps page
        # extraction compatible with the scanned-page image fallback below.
        import fitz

        pdf = fitz.open(str(path))
        try:
            return "\n".join(page.get_text("text") for page in pdf)
        finally:
            pdf.close()
    if suffix == ".docx":
        from docx import Document

        doc = Document(str(path))
        return "\n".join(item.text for item in doc.paragraphs)
    if suffix == ".pptx":
        from pptx import Presentation

        prs = Presentation(str(path))
        lines: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    lines.append(shape.text)
        return "\n".join(lines)
    if suffix == ".zip":
        # A zip of notes: extract text from every supported file inside and concatenate.
        import zipfile

        supported = TEXT_EXTENSIONS | {".pdf", ".docx", ".pptx"}
        parts: list[str] = []
        with zipfile.ZipFile(str(path)) as archive:
            for name in archive.namelist():
                inner = Path(name).suffix.lower()
                if inner not in supported or name.endswith("/"):
                    continue
                with tempfile.NamedTemporaryFile(delete=False, suffix=inner) as handle:
                    handle.write(archive.read(name))
                    inner_path = Path(handle.name)
                try:
                    parts.append(f"[{name}]\n{_extract_from_path(inner_path)}")
                except Exception:  # noqa: BLE001 — skip a single unreadable member
                    pass
                finally:
                    try:
                        os.remove(inner_path)
                    except OSError:
                        pass
        return "\n\n".join(parts)
    return path.read_text(encoding="utf-8", errors="ignore")


def extract_text(source: str | Path) -> str:
    if _is_remote(source):
        response = requests.get(str(source), timeout=60)
        response.raise_for_status()
        suffix = Path(urlparse(str(source)).path).suffix or ".bin"
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as handle:
            handle.write(response.content)
            temp_path = Path(handle.name)
        try:
            return _extract_from_path(temp_path)
        finally:
            try:
                os.remove(temp_path)
            except OSError:
                pass
    return _extract_from_path(Path(source))
